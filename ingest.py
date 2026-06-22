#!/usr/bin/env python3
"""
Ingest a knowledge base into a Chroma vector DB with language-aware chunking.

Walks the data directory and dispatches each file to a format-specific loader:
text/code, PDF, Word, Excel, PowerPoint, and images via OCR. Each format's
dependency is imported lazily inside its loader, so a missing parser disables
only that one file type (with a clear pip hint) instead of breaking the run.
"""
import argparse
import os
from pathlib import Path
from dotenv import load_dotenv
from typing import Callable, Dict, List, Tuple

load_dotenv()

# --- Imports for Splitting and DB ---
from langchain_text_splitters import RecursiveCharacterTextSplitter, Language
from langchain_chroma import Chroma
from langchain_core.documents import Document

DEFAULT_DATA_DIR = os.getenv("DATA_DIR", "./data")
DEFAULT_PERSIST_DIR = os.getenv("CHROMA_PERSIST_DIR", "./chroma_db")
LLM_PROVIDER = os.getenv("LLM_PROVIDER", "google").strip().lower()
HF_EMBEDDING_MODEL = os.getenv("HF_EMBEDDING_MODEL", "all-MiniLM-L6-v2")
OLLAMA_EMBEDDING_MODEL = os.getenv("OLLAMA_EMBEDDING_MODEL", "mxbai-embed-large")
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")

EXTENSION_MAP = {
    ".py": Language.PYTHON,
    ".js": Language.JS,
    ".ts": Language.TS,
    ".java": Language.JAVA,
    ".cpp": Language.CPP,
    ".go": Language.GO,
    ".rb": Language.RUBY,
    ".php": Language.PHP,
    ".rs": Language.RUST,
    ".scala": Language.SCALA,
    ".md": Language.MARKDOWN,
    ".html": Language.HTML,
    # Prose/data formats (.txt, .json, .csv, .yaml, ...) intentionally fall
    # through to the generic character splitter below.
}

def _huggingface_embeddings():
    """Local, zero-key fallback embeddings. Imported lazily."""
    try:
        from langchain_huggingface import HuggingFaceEmbeddings  # current package
    except ImportError:
        # Older installs still expose it under community.
        from langchain_community.embeddings import HuggingFaceEmbeddings  # type: ignore
    print(f"Using HuggingFaceEmbeddings (local): {HF_EMBEDDING_MODEL}")
    return HuggingFaceEmbeddings(model_name=HF_EMBEDDING_MODEL)


def create_embeddings():
    """Return an embedding model that matches the configured provider.

    IMPORTANT: the same embedding model must be used for both ingest and query,
    otherwise vector dimensions won't match. Keep LLM_PROVIDER consistent across
    `ingest.py` and the chat/driver runs.

    - ollama  -> local OllamaEmbeddings (zero keys)
    - openai  -> OpenAIEmbeddings (needs OPENAI_API_KEY)
    - google  -> GoogleGenerativeAIEmbeddings (needs GOOGLE_API_KEY)
    - anything else / missing deps -> local HuggingFace fallback (zero keys)
    """
    # --- LOCAL: Ollama (no key) ---
    if LLM_PROVIDER == "ollama":
        try:
            from langchain_ollama import OllamaEmbeddings
            print(f"Using OllamaEmbeddings (local): {OLLAMA_EMBEDDING_MODEL}")
            return OllamaEmbeddings(model=OLLAMA_EMBEDDING_MODEL, base_url=OLLAMA_BASE_URL)
        except Exception as e:
            print(f"OllamaEmbeddings unavailable ({e}). Falling back to local HuggingFace.")
            return _huggingface_embeddings()

    # --- CLOUD: OpenAI ---
    if LLM_PROVIDER == "openai" and os.getenv("OPENAI_API_KEY"):
        try:
            from langchain_openai import OpenAIEmbeddings
            print("Using OpenAIEmbeddings.")
            return OpenAIEmbeddings()
        except Exception as e:
            print(f"OpenAIEmbeddings unavailable ({e}). Falling back to local HuggingFace.")
            return _huggingface_embeddings()

    # --- CLOUD: Google (Gemini API key, matches the chat provider) ---
    if LLM_PROVIDER == "google" and os.getenv("GOOGLE_API_KEY"):
        try:
            from langchain_google_genai import GoogleGenerativeAIEmbeddings
            print("Using GoogleGenerativeAIEmbeddings.")
            return GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
        except Exception as e:
            print(f"GoogleGenerativeAIEmbeddings unavailable ({e}). Falling back to local HuggingFace.")
            return _huggingface_embeddings()

    # --- LOCAL fallback (no key) ---
    if LLM_PROVIDER in ("google", "openai"):
        print(
            f"WARNING: LLM_PROVIDER='{LLM_PROVIDER}' but no API key was found. "
            "Falling back to local HuggingFace embeddings for ingest. Note: "
            "generation will still require the cloud key, and if you add it "
            "later you must re-run ingest -- the embedding model must match "
            "between ingest and query."
        )
    return _huggingface_embeddings()

def get_splitter_for_file(filename: str, chunk_size: int, chunk_overlap: int):
    """Returns a text splitter optimized for the file's language."""
    _, ext = os.path.splitext(filename)
    lang = EXTENSION_MAP.get(ext.lower())
    
    if lang:
        return RecursiveCharacterTextSplitter.from_language(
            language=lang, chunk_size=chunk_size, chunk_overlap=chunk_overlap
        )
    else:
        # Fallback for unknown extensions (markdown, config files, etc.)
        return RecursiveCharacterTextSplitter(
            chunk_size=chunk_size, chunk_overlap=chunk_overlap
        )


# ---------------------------------------------------------------------------
# Multi-format file loaders
#
# Each format's heavy dependency is imported lazily inside its own loader, so a
# missing parser only disables that one file type (with a clear pip hint)
# instead of breaking the whole ingest -- the same philosophy used for LLM
# providers. Supported inputs:
#   Text/code/data : .py .js .ts .java .cpp .go .rb .php .rs .scala
#                    .txt .md .rst .json .csv .yaml .yml .html .css .sql
#   PDF            : .pdf            (needs `pypdf`)
#   Word           : .docx           (needs `docx2txt`)
#   Excel          : .xlsx .xls       (needs `pandas` + `openpyxl`/`xlrd`)
#   PowerPoint     : .pptx           (needs `python-pptx`)
#   Images (OCR)   : .png .jpg .jpeg .tif .tiff .bmp .gif .webp
#                    (needs `pytesseract` + `Pillow` + the Tesseract engine)
# ---------------------------------------------------------------------------

# Text-like formats are read as UTF-8 with encoding auto-detection.
TEXT_EXTS = {
    ".py", ".js", ".ts", ".java", ".cpp", ".go", ".rb", ".php", ".rs", ".scala",
    ".txt", ".md", ".rst", ".json", ".csv", ".yaml", ".yml", ".html", ".css", ".sql",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}


class LoaderDependencyError(RuntimeError):
    """Raised when a format's optional dependency (or system tool) is missing."""


def _load_text(path: Path) -> List[Document]:
    from langchain_community.document_loaders import TextLoader
    return TextLoader(str(path), encoding="utf-8", autodetect_encoding=True).load()


def _load_pdf(path: Path) -> List[Document]:
    try:
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader(str(path)).load()
    except ImportError as e:
        raise LoaderDependencyError("PDF support needs 'pypdf' (pip install pypdf).") from e


def _load_docx(path: Path) -> List[Document]:
    try:
        from langchain_community.document_loaders import Docx2txtLoader
        return Docx2txtLoader(str(path)).load()
    except ImportError as e:
        raise LoaderDependencyError("Word .docx support needs 'docx2txt' (pip install docx2txt).") from e


def _load_excel(path: Path) -> List[Document]:
    """Read every sheet; emit one Document per sheet as CSV-style text."""
    try:
        import pandas as pd
    except ImportError as e:
        raise LoaderDependencyError("Excel support needs 'pandas' (pip install pandas openpyxl).") from e
    try:
        sheets = pd.read_excel(str(path), sheet_name=None)  # dict: {sheet_name: DataFrame}
    except ImportError as e:
        raise LoaderDependencyError(
            "Excel support needs an engine: 'openpyxl' for .xlsx, 'xlrd' for legacy .xls."
        ) from e
    docs: List[Document] = []
    for name, df in sheets.items():
        text = f"# Sheet: {name}\n{df.to_csv(index=False)}"
        docs.append(Document(page_content=text, metadata={"source": str(path), "sheet": str(name)}))
    return docs


def _load_pptx(path: Path) -> List[Document]:
    """Extract text and tables from each slide; one Document per non-empty slide."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise LoaderDependencyError("PowerPoint .pptx support needs 'python-pptx' (pip install python-pptx).") from e
    prs = Presentation(str(path))
    docs: List[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append(" | ".join(cell.text for cell in row.cells))
        slide_text = "\n".join(lines).strip()
        if slide_text:
            docs.append(Document(
                page_content=f"# Slide {i}\n{slide_text}",
                metadata={"source": str(path), "slide": i},
            ))
    return docs


def _load_image(path: Path) -> List[Document]:
    """OCR an image to text. Returns [] if nothing legible was extracted."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError as e:
        raise LoaderDependencyError(
            "Image OCR needs 'pytesseract' + 'Pillow' (pip install pytesseract pillow) "
            "and the Tesseract engine installed (set TESSERACT_CMD if it's not on PATH)."
        ) from e
    cmd = os.getenv("TESSERACT_CMD")
    if cmd:
        pytesseract.pytesseract.tesseract_cmd = cmd
    text = (pytesseract.image_to_string(Image.open(str(path))) or "").strip()
    if not text:
        return []
    return [Document(page_content=text, metadata={"source": str(path), "type": "image-ocr"})]


# Extension -> loader for binary/structured formats.
_EXT_LOADERS: Dict[str, Callable[[Path], List[Document]]] = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".xlsx": _load_excel,
    ".xls": _load_excel,
    ".pptx": _load_pptx,
}

# Everything we know how to read, for display/help.
SUPPORTED_EXTS = sorted(TEXT_EXTS | IMAGE_EXTS | set(_EXT_LOADERS))


def loader_for(path: Path):
    """Return the loader function for a file, or None if the type is unsupported."""
    ext = path.suffix.lower()
    if ext in _EXT_LOADERS:
        return _EXT_LOADERS[ext]
    if ext in TEXT_EXTS:
        return _load_text
    if ext in IMAGE_EXTS:
        return _load_image
    return None


def load_documents(data_path: Path) -> Tuple[List[Document], List[Tuple[str, str]], int]:
    """Walk the directory and load every supported file via its format loader.

    Returns (documents, skipped, unsupported_count) where `skipped` is a list of
    (path, reason) for files whose parser is missing or that failed to read.
    """
    documents: List[Document] = []
    skipped: List[Tuple[str, str]] = []
    unsupported = 0

    for path in sorted(p for p in data_path.rglob("*") if p.is_file()):
        loader = loader_for(path)
        if loader is None:
            unsupported += 1
            continue
        rel = path.relative_to(data_path)
        try:
            docs = loader(path)
            if docs:
                documents.extend(docs)
                print(f"  + {rel} ({len(docs)} section(s))")
            else:
                print(f"  . {rel} (no extractable text)")
        except Exception as e:
            skipped.append((str(path), str(e)))
            print(f"  ! skipped {rel}: {e}")

    return documents, skipped, unsupported


def ingest_folder(data_dir: str, persist_dir: str, chunk_size: int = 2000, chunk_overlap: int = 200):
    data_path = Path(data_dir)
    if not data_path.exists():
        raise FileNotFoundError(f"{data_dir} not found.")

    print(f"Scanning {data_dir} for supported files...")
    print(f"Supported types: {', '.join(SUPPORTED_EXTS)}\n")
    documents, skipped, unsupported = load_documents(data_path)

    if skipped:
        print(f"\n{len(skipped)} file(s) skipped (missing parser or read error). "
              "Install the optional dependency named in each message to enable them.")
    if unsupported:
        print(f"{unsupported} file(s) ignored (unsupported extension).")

    if not documents:
        print("No documents loaded. Nothing to ingest.")
        return None

    # Language-aware chunking, per source file.
    print(f"\nSplitting {len(documents)} document section(s) with language-aware logic...")
    final_chunks: List[Document] = []
    for doc in documents:
        source_file = doc.metadata.get("source", "")
        splitter = get_splitter_for_file(source_file, chunk_size, chunk_overlap)
        final_chunks.extend(splitter.split_documents([doc]))

    embeddings = create_embeddings()
    print(f"Ingesting {len(final_chunks)} chunks...")
    vectordb = Chroma.from_documents(
        documents=final_chunks,
        embedding=embeddings,
        persist_directory=persist_dir,
    )
    print(f"Done. Ingested {len(final_chunks)} chunks into Chroma at {persist_dir}")
    return vectordb


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--data-dir", default=DEFAULT_DATA_DIR,
                        help="Directory of documents to ingest (text/code, PDF, Word, Excel, PowerPoint, images)")
    parser.add_argument("--persist-dir", default=DEFAULT_PERSIST_DIR, help="Chroma persist directory")
    parser.add_argument("--chunk-size", type=int, default=2000)
    parser.add_argument("--chunk-overlap", type=int, default=200)
    args = parser.parse_args()
    ingest_folder(args.data_dir, args.persist_dir, args.chunk_size, args.chunk_overlap)
